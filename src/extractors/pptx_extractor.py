from pptx import Presentation

class PptxExtractor:

    def extract(self, file_path):

        prs = Presentation(file_path)

        text = []

        for slide in prs.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    if shape.text.strip():
                        text.append(shape.text)

        return "\n".join(text)